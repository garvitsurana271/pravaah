"""Generates Pravaah-Deck.pptx — the FAR AWAY 2026 submission deck."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SHOTS = os.path.join(ROOT, "docs", "screenshots")

BG = RGBColor(0x02, 0x06, 0x17)
PANEL = RGBColor(0x10, 0x18, 0x2c)
GREEN = RGBColor(0x22, 0xd3, 0x7a)
BLUE = RGBColor(0x3a, 0xa0, 0xff)
AMBER = RGBColor(0xff, 0xb0, 0x2e)
RED = RGBColor(0xff, 0x4d, 0x4d)
INK = RGBColor(0xe9, 0xf0, 0xfa)
MUTED = RGBColor(0x9a, 0xb0, 0xcf)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s


def bar(s, color):
    b = s.shapes.add_shape(1, 0, 0, Inches(0.14), SH)
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()


def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    return tf


def para(tf, text, size, color, bold=False, font="Segoe UI", space=8, align=PP_ALIGN.LEFT, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.space_after = Pt(space)
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
    r.font.color.rgb = color
    return p


def kicker(s, text, color=GREEN):
    tf = tb(s, Inches(0.7), Inches(0.5), Inches(12), Inches(0.5))
    para(tf, text.upper(), 13, color, bold=True, font="Consolas", first=True)


def title(s, text, color=INK, y=0.95, size=40):
    tf = tb(s, Inches(0.7), Inches(y), Inches(12), Inches(1.3))
    para(tf, text, size, color, bold=True, first=True)


def bullets(s, items, x=0.7, y=2.2, w=6.2, size=18):
    tf = tb(s, Inches(x), Inches(y), Inches(w), Inches(4.6))
    for i, (txt, col) in enumerate(items):
        p = para(tf, txt, size, col, space=14, first=(i == 0))
        p.level = 0


def pic(s, name, x, y, w):
    path = os.path.join(SHOTS, name)
    if os.path.exists(path):
        s.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))


def footer(s, text):
    tf = tb(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4))
    para(tf, text, 11, MUTED, font="Consolas", first=True)


# 1 — TITLE
s = slide(); bar(s, GREEN)
tf = tb(s, Inches(0.9), Inches(2.1), Inches(11.5), Inches(3.2))
para(tf, "FAR AWAY · 2026   ·   RAILWAYS", 15, BLUE, bold=True, font="Consolas", first=True)
para(tf, "PRAVAAH  प्रवाह", 60, INK, bold=True, space=4)
para(tf, "The Glass-Box Dispatcher", 30, GREEN, bold=True, space=16)
para(tf, "An open, explainable, real-time AI train re-dispatch copilot for the Indian Railways section controller.", 19, MUTED)
footer(s, "github.com/garvitsurana271/pravaah   ·   maps to SIH25022   ·   Garvit Surana & team")

# 2 — THE TRAGEDY
s = slide(); bar(s, RED); kicker(s, "the problem · part 1", RED); title(s, "The next Balasore shouldn't depend on one tired controller's memory")
bullets(s, [
    ("2 June 2023 — a single mis-wired signalling circuit sent the Coromandel Express onto an occupied loop line at 128 km/h near Bahanaga Bazar, Balasore.", INK),
    ("~300 dead. The line had no Kavach.", RED),
    ("June 2024 — a goods train ran a red signal (SPAD) and killed ten more on the Kanchanjunga Express. Again, no Kavach.", MUTED),
], w=11.5, size=20)

# 3 — THE DAILY PROBLEM
s = slide(); bar(s, AMBER); kicker(s, "the problem · part 2", AMBER); title(s, "Indian Railways dispatches every train by hand")
bullets(s, [
    ("A human section controller decides every crossing and precedence call from experience.", INK),
    ("Control-room software (COA, ICMS, NTES) charts and monitors — it does not optimise.", MUTED),
    ("Punctuality fell 94% (FY20-21) → 73.6% (FY23-24).", AMBER),
    ("Over 80% of the busiest corridors run above capacity.", AMBER),
], w=11.5, size=20)

# 4 — THE GAP (novelty)
s = slide(); bar(s, GREEN); kicker(s, "the gap = our novelty"); title(s, "Proprietary ₹1000-crore black boxes — or nothing")
bullets(s, [
    ("Hitachi / Siemens / Alstom sell closed, crore-scale traffic-management systems — and even those stay advisory, because controllers don't trust black boxes.", INK),
    ("There is NO open, India-specific, explainable tool for this.", GREEN),
    ("Pravaah's novelty is the stack: open + explainable + India + real-time recovery — not the underlying operations-research.", MUTED),
    ("We deliberately avoided RL-for-dispatching: the Flatland challenge showed classical OR beats it, and it has never deployed.", MUTED),
], w=11.5, size=19)

# 5 — WHAT IT IS (control room)
s = slide(); bar(s, BLUE); kicker(s, "the product"); title(s, "A live control room for one congested section", size=34)
pic(s, "control-room.png", 0.7, 2.05, 12.0)
footer(s, "Two brains run side-by-side on the same scenario: the AI optimizer and the manual first-come-first-served baseline.")

# 6 — THROUGHPUT A/B
s = slide(); bar(s, GREEN); kicker(s, "result · throughput"); title(s, "Flip the optimizer off — watch the delay jump", size=34)
bullets(s, [
    ("Same eight trains, same track. The optimizer minimises Σ (delay × priority-weight).", INK),
    ("–20% weighted delay vs. the manual FCFS baseline.", GREEN),
    ("It doesn't move more trains — it protects the premier ones, spending a goods rake's minutes to save a Superfast's.", MUTED),
    ("13 conflicts auto-resolved. 0 unsafe states. (measured by the test suite)", MUTED),
], w=11.5, size=20)

# 7 — GLASS BOX
s = slide(); bar(s, GREEN); kicker(s, "the differentiator"); title(s, "It explains every call — in plain language", size=32)
bullets(s, [
    ("Every decision is the solver's own trace, rendered as auditable reasoning — never invented.", INK),
    ("“Held 18045 (Express, ×3) so the Superfast could cross — reversing would cost 1.7× more.”", GREEN),
    ("This directly answers why AI dispatching isn't adopted: black-box distrust.", MUTED),
], x=0.7, y=2.0, w=5.7, size=17)
pic(s, "glass-box.png", 6.7, 1.7, 6.0)

# 8 — COPILOT
s = slide(); bar(s, BLUE); kicker(s, "interrogate it"); title(s, "Ask the dispatcher anything — offline", size=32)
bullets(s, [
    ("“Why is 12841 held?”  “Is the section safe?”  “What's the cost of overriding?”", INK),
    ("Grounded in the same decision trace. No network call — the demo cannot break.", MUTED),
    ("The controller stays in command; the machine just shows the cost of each choice.", MUTED),
], x=0.7, y=2.0, w=5.7, size=17)
pic(s, "glass-box.png", 6.7, 1.7, 6.0)

# 9 — SAFETY FLOOR
s = slide(); bar(s, RED); kicker(s, "the safety floor", RED); title(s, "Unsafe states are impossible by construction", size=32)
bullets(s, [
    ("A separate interlocking layer enforces one-train-per-block and single-line direction-locks — independent of any AI policy.", INK),
    ("Re-stage the Balasore failure: a route mis-set toward an occupied line is REFUSED.", RED),
    ("Safety is a hard invariant; the AI only optimises within the safe envelope.", MUTED),
], x=0.7, y=2.0, w=5.6, size=16)
pic(s, "safety-hold.png", 6.6, 1.7, 6.1)

# 10 — ARCHITECTURE
s = slide(); bar(s, BLUE); kicker(s, "how it works"); title(s, "Two separated layers")
bullets(s, [
    ("INTERLOCKING (safety) — one train per block · single-line direction-lock · loop limits. Unsafe = impossible.", GREEN),
    ("DISPATCHER (policy) — OPTIMIZER vs FCFS; minimises weighted delay; emits a decision trace.", BLUE),
    ("GLASS-BOX ENGINE — turns that trace into plain language + answers questions.", INK),
    ("Deterministic, tick-based, pure TypeScript. 100% client-side — no backend, no API keys.", MUTED),
], w=11.5, size=19)

# 11 — MAPS TO SIH25022
s = slide(); bar(s, GREEN); kicker(s, "validated demand"); title(s, "We built the Ministry's own flagship problem")
bullets(s, [
    ("SIH25022 (Smart India Hackathon, Ministry of Railways):", MUTED),
    ("“Maximizing section throughput through AI-powered, real-time train traffic control.”", GREEN),
    ("Complementary to Kavach: Kavach is the safety-certified anti-collision floor; Pravaah is the throughput-and-explainability layer above it.", INK),
], w=11.5, size=20)

# 12 — TECH + RIGOR
s = slide(); bar(s, BLUE); kicker(s, "engineering"); title(s, "Built to be inspected")
bullets(s, [
    ("React 18 + TypeScript + Vite + Tailwind + SVG. Bundled fonts → fully offline.", INK),
    ("12 passing tests: safety invariants, deadlock-freedom, optimizer-beats-FCFS, the safety refusal.", GREEN),
    ("Real Indian station names & coordinates from the open datameet/railways dataset (CC0).", MUTED),
    ("One command to run: npm install && npm run dev.", MUTED),
], w=11.5, size=19)
pic(s, "disruption.png", 7.4, 4.4, 5.4)

# 13 — CLOSE
s = slide(); bar(s, GREEN)
tf = tb(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(3))
para(tf, "Pravaah · प्रवाह", 44, INK, bold=True, first=True)
para(tf, "Open. Explainable. Safe by construction. The control-room brain Indian Railways doesn't have yet.", 22, GREEN, bold=True, space=18)
para(tf, "Honest scope: a decision-support prototype — complementary to, not a replacement for, India's safety-certified train control.", 15, MUTED)
footer(s, "github.com/garvitsurana271/pravaah   ·   FAR AWAY 2026   ·   Railways   ·   SIH25022")

out = os.path.join(ROOT, "Pravaah-Deck.pptx")
prs.save(out)
print("saved", out, "·", len(prs.slides._sldIdLst), "slides")
