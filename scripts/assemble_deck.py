"""Assembles Pravaah-Deck.pptx from the slide images rendered from
public/slides.html (each a full-bleed 1920x1080 PNG). This is how the deck is
built: design the slides in HTML/CSS, render them, then package as PPTX, so the
deck looks designed rather than auto-generated."""
import os
from pptx import Presentation
from pptx.util import Inches

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

for i in range(1, 11):
    path = os.path.join(ROOT, f"slide-preview-{i}.png")
    if not os.path.exists(path):
        raise SystemExit(f"missing {path} (render public/slides.html first)")
    s = prs.slides.add_slide(BLANK)
    s.shapes.add_picture(path, 0, 0, Inches(13.333), Inches(7.5))

out = os.path.join(ROOT, "Pravaah-Deck.pptx")
prs.save(out)
print("saved", out, "·", len(prs.slides._sldIdLst), "slides ·", os.path.getsize(out), "bytes")
